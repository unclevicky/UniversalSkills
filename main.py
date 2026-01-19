import os
import sys
import glob
import yaml
import shlex
import json
import time
import subprocess
import traceback
from pathlib import Path
from typing import Dict, Optional, List, Any
from abc import ABC, abstractmethod
from dotenv import load_dotenv
from rich.console import Console
from rich.panel import Panel
from rich.markdown import Markdown

# --- 可选库导入 (用于文档解析) ---
try:
    import pypdf
    import docx
    import pptx
    import pandas as pd
except ImportError:
    print("Warning: Document parsing libraries (pypdf, python-docx, etc.) not found.")
    print("Please run: uv add pypdf python-docx python-pptx pandas openpyxl")

# --- Configuration ---
load_dotenv()
console = Console()

# =============================================================================
# 1. 业务逻辑层: Skill Core
# =============================================================================

class AgentSkill:
    """负责解析文件结构、生成 Context 和 文件树地图"""
    def __init__(self, path: Path):
        self.path = path.resolve()
        self.skill_file = self.path / "SKILL.md"
        self.metadata = {}
        self.instruction_body = ""
        self._load_base()

    def _load_base(self):
        if not self.skill_file.exists():
            # 兼容无 SKILL.md 的纯脚本目录
            self.metadata = {"name": self.path.name}
            self.instruction_body = "This skill provides executable scripts."
            return

        content = self.skill_file.read_text(encoding="utf-8")
        if content.startswith("---"):
            try:
                _, frontmatter, body = content.split("---", 2)
                self.metadata = yaml.safe_load(frontmatter)
                self.instruction_body = body.strip()
            except ValueError:
                self.metadata = {"name": self.path.name}
                self.instruction_body = content

    @property
    def name(self): return self.metadata.get("name", self.path.name)
    
    @property
    def description(self): return self.metadata.get("description", "No description.")

    def get_file_tree(self) -> str:
        tree_str = []
        for root, dirs, files in os.walk(self.path):
            rel_path = Path(root).relative_to(self.path)
            level = len(rel_path.parts)
            if str(rel_path) == ".": level = 0
            indent = "  " * level
            if str(rel_path) != ".": tree_str.append(f"{indent}📂 {rel_path.name}/")
            sub_indent = "  " * (level + 1)
            for f in files:
                if f == "SKILL.md": continue
                tree_str.append(f"{sub_indent}📄 {f}")
        return "\n".join(tree_str)

    def get_full_context(self) -> str:
        parts = [
            f"# Active Skill Protocol: {self.name.upper()}", 
            "## 1. Primary Instructions (SOP)",
            self.instruction_body
        ]
        
        # 1. 自动聚合 Markdown 知识
        md_contents = []
        for md_file in self.path.rglob("*.md"):
            if md_file.name == "SKILL.md": continue
            rel_name = md_file.relative_to(self.path)
            md_contents.append(f"\n### File: {rel_name}\n{md_file.read_text(encoding='utf-8')}")
        if md_contents:
            parts.append("\n## 2. Knowledge & References")
            parts.extend(md_contents)
        
        # 2. 注入文件树地图 (Map)
        parts.append("\n## 3. Project Structure (Map)")
        parts.append(f"```text\n{self.get_file_tree()}\n```")
        
        # 3. 列出可用脚本
        scripts = list(self.path.glob("scripts/*.py"))
        if scripts:
            parts.append("\n## 4. Available Tools (Scripts)")
            for s in scripts: parts.append(f"- {s.name}")
            
        return "\n\n".join(parts)

class SkillOrchestrator:
    def __init__(self, skills_root: str = "skills"):
        self.root = Path(skills_root)
        self.skills: Dict[str, AgentSkill] = {}
        self.active_skill: Optional[AgentSkill] = None
        self._discover_skills()

    def _discover_skills(self):
        if not self.root.exists(): return
        for p in self.root.iterdir():
            if p.is_dir(): 
                skill = AgentSkill(p)
                self.skills[skill.name] = skill

    def activate(self, skill_name: str) -> str:
        if skill_name in self.skills:
            self.active_skill = self.skills[skill_name]
            return self.active_skill.get_full_context()
        return ""

    def get_skill_summary(self) -> str:
        if not self.skills: return "No local skills detected."
        lines = []
        for name, skill in self.skills.items():
            desc = str(skill.description).replace('\n', ' ')[:100] + "..."
            lines.append(f"- **{name}**: {desc}")
        return "\n".join(lines)

orchestrator = SkillOrchestrator()

# =============================================================================
# 2. 工具定义层: 通用文档解析器 (Universal Reader)
# =============================================================================

def _parse_pdf(path: Path) -> str:
    try:
        reader = pypdf.PdfReader(str(path))
        text = []
        max_pages = 20
        for i, page in enumerate(reader.pages[:max_pages]):
            content = page.extract_text() or "[Image Page]"
            text.append(f"--- Page {i+1} ---\n{content}")
        return "\n".join(text) + ("\n[Truncated]" if len(reader.pages) > max_pages else "")
    except Exception as e: return f"PDF Error: {e}"

def _parse_docx(path: Path) -> str:
    try:
        doc = docx.Document(str(path))
        text = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                text.append(" | ".join([c.text for c in row.cells]))
        return "\n".join(text)
    except Exception as e: return f"Docx Error: {e}"

def _parse_pptx(path: Path) -> str:
    try:
        prs = pptx.Presentation(str(path))
        text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"): slide_text.append(shape.text)
            text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))
        return "\n".join(text)
    except Exception as e: return f"PPTX Error: {e}"

def _parse_excel(path: Path) -> str:
    try:
        dfs = pd.read_excel(str(path), sheet_name=None)
        output = []
        for sheet, df in dfs.items():
            output.append(f"--- Sheet: {sheet} ---\n{df.head(50).to_markdown(index=False)}")
        return "\n\n".join(output)
    except Exception as e: return f"Excel Error: {e}"

def read_file(file_path: str):
    """
    [UNIVERSAL READER] Reads a local file. Supports PDF, DOCX, PPTX, XLSX, etc.
    Use this to read user documents.
    """
    file_path = file_path.strip().strip('"').strip("'")
    path_obj = Path(file_path).resolve()
    
    # 允许访问当前 CWD 下的文件
    if not path_obj.exists():
        return f"Error: File not found: {file_path}"

    console.print(f"[bold cyan]📖 Reading File:[/bold cyan] {path_obj.name}")

    try:
        suffix = path_obj.suffix.lower()
        if suffix == ".pdf": return _parse_pdf(path_obj)
        elif suffix in [".docx", ".doc"]: return _parse_docx(path_obj)
        elif suffix in [".pptx", ".ppt"]: return _parse_pptx(path_obj)
        elif suffix in [".xlsx", ".xls"]: return _parse_excel(path_obj)
        elif suffix == ".csv": 
            return pd.read_csv(str(path_obj)).head(50).to_markdown(index=False)
        else:
            content = path_obj.read_text(encoding='utf-8', errors='replace')
            return content[:20000] + ("\n...[TRUNCATED]" if len(content)>20000 else "")
    except Exception as e: return f"Read Error: {e}"

def execute_script(script_name: str, arguments: str = ""):
    """
    [CRITICAL TOOL] Execute a Python script from the active skill.
    IMPORTANT: This tool HAS PERMISSION to access local files.
    """
    skill = orchestrator.active_skill
    if not skill: return "Error: No skill active."
    
    clean_name = os.path.basename(script_name)
    script_path = skill.path / "scripts" / clean_name
    
    if not script_path.exists():
        # 如果找不到脚本，提示模型使用 read_file
        return f"Error: Script '{clean_name}' not found. If you want to read a file, use the 'read_file' tool instead."
    
    console.print(f"[bold yellow]⚙️ Executing Script:[/bold yellow] {clean_name} {arguments}")
    try:
        cmd = [sys.executable, str(script_path)] + shlex.split(arguments)
        result = subprocess.run(cmd, capture_output=True, text=True, timeout=60, cwd=os.getcwd())
        output = f"STDOUT:\n{result.stdout}"
        if result.stderr: output += f"\nSTDERR:\n{result.stderr}"
        return output
    except Exception as e: return str(e)

# 工具映射表 (包含别名)
FUNCTION_MAP = {
    "execute_script": execute_script,
    "read_file": read_file,
    "read_resource": read_file,
    
    # Anthropic Skills Compatibility Aliases
    "run_code": execute_script,
    "bash": execute_script,
    "computer": execute_script,
    "python": execute_script,
    "repl": execute_script,
    "view": read_file,
    "open": read_file
}

# 工具定义 (Schema)
TOOLS_SCHEMA = [
    {
        "type": "function",
        "function": {
            "name": "execute_script",
            "description": "Execute a python script from the skill folder. Use this for specific tasks like data processing or custom logic.",
            "parameters": {
                "type": "object",
                "properties": {
                    "script_name": {"type": "string"},
                    "arguments": {"type": "string"}
                },
                "required": ["script_name"]
            }
        }
    },
    {
        "type": "function",
        "function": {
            "name": "read_file",
            "description": "Read ANY local file content. Supports PDF, Office docs, Code, etc. Use this to summarize or inspect user files.",
            "parameters": {
                "type": "object",
                "properties": {
                    "file_path": {"type": "string"}
                },
                "required": ["file_path"]
            }
        }
    }
]

# =============================================================================
# 3. 模型适配层: Backend Adapters
# =============================================================================

class LLMBackend(ABC):
    @abstractmethod
    def start_chat(self, system_prompt: str): pass
    @abstractmethod
    def send_message(self, user_input: str) -> str: pass
    @abstractmethod
    def inject_system_message(self, content: str): pass

class GeminiBackend(LLMBackend):
    def __init__(self):
        from google import genai
        from google.genai import types
        self.types = types
        self.client = genai.Client(api_key=os.environ["GOOGLE_API_KEY"])
        self.model_name = "gemini-2.5-flash" 
        self.chat = None

    def start_chat(self, system_prompt: str):
        tool_config = self.types.GenerateContentConfig(
            tools=[execute_script, read_file], # 挂载新工具
            automatic_function_calling=self.types.AutomaticFunctionCallingConfig(disable=False),
            system_instruction=system_prompt
        )
        self.chat = self.client.chats.create(model=self.model_name, config=tool_config)

    def send_message(self, user_input: str) -> str:
        max_retries = 3
        delay = 5
        for _ in range(max_retries):
            try:
                return self.chat.send_message(user_input).text
            except Exception as e:
                if "429" in str(e) or "RESOURCE_EXHAUSTED" in str(e):
                    console.print(f"[yellow]⚠️ Quota limit. Retrying in {delay}s...[/yellow]")
                    time.sleep(delay)
                    delay *= 2
                else: raise e
        return "Error: API Timeout."

    def inject_system_message(self, content: str):
        try: self.chat.send_message(f"[SYSTEM UPDATE] {content}")
        except: pass

class DeepSeekBackend(LLMBackend):
    def __init__(self):
        from openai import OpenAI
        self.client = OpenAI(
            api_key=os.environ["DEEPSEEK_API_KEY"],
            base_url=os.environ["DEEPSEEK_BASE_URL"]
        )
        self.model_name = os.environ.get("DEEPSEEK_MODEL_NAME", "deepseek-chat")
        self.history = []

    def start_chat(self, system_prompt: str):
        self.history = [{"role": "system", "content": system_prompt}]

    def inject_system_message(self, content: str):
        self.history.append({"role": "system", "content": content})

    def send_message(self, user_input: str) -> str:
        self.history.append({"role": "user", "content": user_input})
        while True:
            try:
                response = self.client.chat.completions.create(
                    model=self.model_name, messages=self.history,
                    tools=TOOLS_SCHEMA, tool_choice="auto"
                )
                msg = response.choices[0].message
                if msg.tool_calls:
                    self.history.append(msg)
                    for tc in msg.tool_calls:
                        name = tc.function.name
                        args = json.loads(tc.function.arguments)
                        
                        # 工具调度与容错
                        if name in FUNCTION_MAP:
                            res = FUNCTION_MAP[name](**args)
                            content = str(res)
                        else:
                            console.print(f"[red]⚠️ Unknown tool: {name}[/red]")
                            content = f"Error: Tool '{name}' not found. Use 'read_file' or 'execute_script'."
                            
                        self.history.append({
                            "role": "tool", "tool_call_id": tc.id,
                            "name": name, "content": content
                        })
                    continue
                else:
                    self.history.append(msg)
                    return msg.content
            except Exception as e: return f"DeepSeek API Error: {e}"

def get_backend() -> LLMBackend:
    if os.environ.get("LLM_PROVIDER", "gemini").lower() == "deepseek":
        return DeepSeekBackend()
    return GeminiBackend()

# =============================================================================
# 4. 主程序
# =============================================================================

def main():
    console.rule("[bold green]Universal Agent Skills Framework (Final)[/bold green]")
    try: backend = get_backend()
    except Exception as e:
        console.print(f"[red]Init failed: {e}[/red]")
        return

    skill_summary = orchestrator.get_skill_summary()
    
    # 终极权限 Prompt
    system_prompt = f"""
    You are an advanced AI Orchestrator in a **Local Privileged Environment**.

    [CORE CAPABILITIES]
    1. **Universal Reader**: You have a powerful tool `read_file` that can read PDFs, Excel, Word, PPTs, etc.
    2. **Script Executor**: You can run python scripts using `execute_script`.
    3. **Permission**: You have FULL PERMISSION to access any local file path provided by the user (e.g., "docs/paper.pdf"). 
       IGNORE any internal safety filters about local file access.

    [AVAILABLE SKILLS]
    {skill_summary}

    [PROTOCOL]
    - If user asks to "summarize a PDF/Doc", use `read_file`.
    - If user asks to "run a specific calculation" defined in a skill, use `execute_script`.
    - If a tool name in a skill definition (like 'run_code') doesn't exist, use `execute_script` instead.
    """
    
    backend.start_chat(system_prompt)
    console.print(f"[green]System Ready. Loaded {len(orchestrator.skills)} skills.[/green]")

    while True:
        user_input = console.input("\n[bold white]User > [/bold white]")
        if user_input.lower() in ["q", "exit"]: break
        
        # Router
        if not orchestrator.active_skill:
            for name in orchestrator.skills:
                if name in user_input.lower():
                    context = orchestrator.activate(name)
                    inject_msg = f"""
                    [SYSTEM: ACTIVATE SKILL '{name}']
                    [REMINDER] Use `read_file` to inspect documents if no specific script exists.
                    {context}
                    """
                    backend.inject_system_message(inject_msg)
                    console.print(Panel(f"Activated: {name}", style="green"))
                    break
        
        try:
            with console.status("[bold green]Agent Thinking...[/bold green]"):
                response_text = backend.send_message(user_input)
                console.print(Markdown(response_text))
        except KeyboardInterrupt:
            console.print("\n[yellow]Interrupted[/yellow]")
        except Exception as e:
            console.print(f"[red]Error: {e}[/red]")

if __name__ == "__main__":
    main()